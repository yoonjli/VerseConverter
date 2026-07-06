import sys
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Usage ──────────────────────────────────────────────────────────────
# python make_bilingual_pptx.py korean.txt english.txt [output.pptx]
# ───────────────────────────────────────────────────────────────────────

def parse_verses(filepath):
    """
    Parses a scripture text file into a list of (reference, text) tuples.

    Supports two layouts:

    Layout A – reference AFTER text (as in the sample documents):
        베드로와 요한이 아직도 ...
        사도행전 4:1

    Layout B – reference BEFORE text (common alternate format):
        사도행전 4:1
        베드로와 요한이 아직도 ...

    A line is treated as a reference if it matches the pattern:
        <BookName> <chapter>:<verse>   (e.g. "사도행전 4:1" or "Acts 4:1")
    """
    ref_pattern = re.compile(r'^.{2,30}\s+\d+:\d+\s*$')

    with open(filepath, encoding='utf-8') as f:
        raw_lines = [l.rstrip() for l in f.readlines()]

    # Remove blank lines for easier processing
    lines = [l for l in raw_lines if l.strip()]

    verses = []
    i = 0
    while i < len(lines):
        line = lines[i]
        is_ref = bool(ref_pattern.match(line))

        # Layout A: text line(s) followed by a reference
        if not is_ref and i + 1 < len(lines) and ref_pattern.match(lines[i + 1]):
            text_parts = [line]
            # Gather any continuation lines before the reference
            j = i + 1
            while j < len(lines) and not ref_pattern.match(lines[j]):
                text_parts.append(lines[j])
                j += 1
            ref = lines[j].strip() if j < len(lines) else ""
            verses.append((ref, " ".join(text_parts).strip()))
            i = j + 1

        # Layout B: reference line followed by text
        elif is_ref and i + 1 < len(lines) and not ref_pattern.match(lines[i + 1]):
            ref = line.strip()
            text_parts = []
            j = i + 1
            while j < len(lines) and not ref_pattern.match(lines[j]):
                text_parts.append(lines[j])
                j += 1
            verses.append((ref, " ".join(text_parts).strip()))
            i = j

        else:
            i += 1  # skip unmatched lines

    return verses


# ── Colors ─────────────────────────────────────────────────────────────
BLACK      = RGBColor(0x00, 0x00, 0x00)
GOLD       = RGBColor(0xD4, 0xAF, 0x37)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GOLD = RGBColor(0xF0, 0xD9, 0x80)
DIVIDER    = RGBColor(0xD4, 0xAF, 0x37)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def add_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, text, left, top, width, height,
                font_size=22, bold=False, color=WHITE,
                align=PP_ALIGN.CENTER, font_name=None):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name or (
        "Malgun Gothic" if any(ord(c) > 127 for c in text) else "Georgia"
    )


def make_slide(prs, ref_ko, korean, ref_en, english):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_bg(slide, BLACK)
    W, H = SLIDE_W, SLIDE_H
    m = Inches(0.65)

    # Top gold bar
    b = slide.shapes.add_shape(1, Inches(0), Inches(0), W, Inches(0.07))
    b.fill.solid(); b.fill.fore_color.rgb = GOLD; b.line.fill.background()

    # Combined reference label centered at top
    ref_label = f"{ref_ko}  /  {ref_en}"
    add_textbox(slide, ref_label, m, Inches(0.15), W - m * 2, Inches(0.55),
                font_size=30, bold=True, color=GOLD, align=PP_ALIGN.CENTER,
                font_name="Malgun Gothic")

    # Korean text (top half)
    add_textbox(slide, korean, m, Inches(0.85), W - m * 2, Inches(2.7),
                font_size=36, color=WHITE, align=PP_ALIGN.CENTER,
                font_name="Malgun Gothic")

    # Divider
    div = slide.shapes.add_shape(1, Inches(1.5), Inches(3.7), W - Inches(3.0), Inches(0.04))
    div.fill.solid(); div.fill.fore_color.rgb = DIVIDER; div.line.fill.background()

    # English text (bottom half)
    add_textbox(slide, english, m, Inches(3.85), W - m * 2, Inches(2.7),
                font_size=36, color=LIGHT_GOLD, align=PP_ALIGN.CENTER,
                font_name="Georgia")

    # Bottom gold bar
    b2 = slide.shapes.add_shape(1, Inches(0), H - Inches(0.07), W, Inches(0.07))
    b2.fill.solid(); b2.fill.fore_color.rgb = GOLD; b2.line.fill.background()


def main():
    if len(sys.argv) < 3:
        print("Usage: python make_bilingual_pptx.py korean.txt english.txt [output.pptx]")
        sys.exit(1)

    ko_file  = sys.argv[1]
    en_file  = sys.argv[2]
    out_file = sys.argv[3] if len(sys.argv) > 3 else "bilingual_scripture.pptx"

    print(f"Parsing {ko_file} ...")
    ko_verses = parse_verses(ko_file)
    print(f"  → {len(ko_verses)} verses found")

    print(f"Parsing {en_file} ...")
    en_verses = parse_verses(en_file)
    print(f"  → {len(en_verses)} verses found")

    if len(ko_verses) != len(en_verses):
        print(f"⚠  Warning: verse counts differ ({len(ko_verses)} KO vs {len(en_verses)} EN). "
              f"Using the shorter count.")

    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    count = min(len(ko_verses), len(en_verses))
    for i in range(count):
        ref_ko, text_ko = ko_verses[i]
        ref_en, text_en = en_verses[i]
        make_slide(prs, ref_ko, text_ko, ref_en, text_en)

    prs.save(out_file)
    print(f"✓ Saved {count} slides → {out_file}")


if __name__ == "__main__":
    main()