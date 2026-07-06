from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

verses = [
    ("사도행전 4:1 / Acts 4:1", "베드로와 요한이 아직도 사람들에게 말하고 있는데, 제사장들과 성전 경비대장과 사두개파 사람들이 몰려왔다.", "The priests and the captain of the temple guard and the Sadducees came up to Peter and John while they were speaking to the people."),
    ("사도행전 4:2 / Acts 4:2", "그들은 사도들이 백성을 가르치는 것과, 예수의 부활을 내세워서 죽은 사람들의 부활을 선전하고 있는 것에 격분해서,", "They were greatly disturbed because the apostles were teaching the people, proclaiming in Jesus the resurrection of the dead."),
    ("사도행전 4:3 / Acts 4:3", "사도들을 붙잡았으나, 날이 이미 저물었으므로 다음 날까지 가두어 두었다.", "They seized Peter and John and, because it was evening, they put them in jail until the next day."),
    ("사도행전 4:4 / Acts 4:4", "그런데 사도들의 말을 들은 사람들 가운데서 믿는 사람이 많으니, 남자 어른의 수가 약 오천 명이나 되었다.", "But many who heard the message believed; so the number of men who believed grew to about five thousand."),
    ("사도행전 4:5 / Acts 4:5", "이튿날 유대의 지도자들과 장로들과 율법학자들이 예루살렘에 모였는데,", "The next day the rulers, the elders and the teachers of the law met in Jerusalem."),
    ("사도행전 4:6 / Acts 4:6", "대제사장 안나스를 비롯해서, 가야바와 요한과 알렉산더와 그 밖에 대제사장의 가문에 속한 사람들이 모두 참석하였다.", "Annas the high priest was there, and so were Caiaphas, John, Alexander and others of the high priest's family."),
    ("사도행전 4:7 / Acts 4:7", "그들은 사도들을 가운데에 세워 놓고서 물었다. \"그대들은 대체 무슨 권세와 누구의 이름으로 이런 일을 하였소?\"", "They had Peter and John brought before them and began to question them: \"By what power or what name did you do this?\""),
    ("사도행전 4:8 / Acts 4:8", "그때에 베드로가 성령이 충만하여 그들에게 말하였다. \"백성의 지도자들과 장로 여러분,", "Then Peter, filled with the Holy Spirit, said to them: \"Rulers and elders of the people!"),
    ("사도행전 4:9 / Acts 4:9", "우리가 오늘 신문을 받는 것이, 병자에게 행한 착한 일과 또 그가 누구의 힘으로 낫게 되었느냐 하는 문제 때문이라면,", "If we are being called to account today for an act of kindness shown to a man who was lame and are being asked how he was healed,"),
    ("사도행전 4:10 / Acts 4:10", "여러분 모두와 모든 이스라엘 백성은 이것을 알아야 합니다. 이 사람이 성한 몸으로 여러분 앞에 서게 된 것은, 여러분이 십자가에 못 박아 죽였으나 하나님이 죽은 사람들 가운데서 살리신 나사렛 예수 그리스도의 이름을 힘입어서 된 것입니다.", "then know this, you and all the people of Israel: It is by the name of Jesus Christ of Nazareth, whom you crucified but whom God raised from the dead, that this man stands before you healed."),
    ("사도행전 4:11 / Acts 4:11", "이 예수는 '너희들 집 짓는 사람들에게는 버림받은 돌이지만, 집 모퉁이의 머릿돌이 되신 분'입니다.", "Jesus is 'the stone you builders rejected, which has become the cornerstone.'"),
    ("사도행전 4:12 / Acts 4:12", "이 예수 밖에는, 다른 아무에게도 구원은 없습니다. 사람들에게 주신 이름 가운데 우리가 의지하여 구원을 얻어야 할 이름은, 하늘 아래에 이 이름 밖에 다른 이름이 없습니다.", "Salvation is found in no one else, for there is no other name under heaven given to mankind by which we must be saved."),
    ("사도행전 4:13 / Acts 4:13", "그들은 베드로와 요한이 본래 배운 것이 없는 보잘것없는 사람인 줄 알았는데, 이렇게 담대하게 말하는 것을 보고 놀랐다. 그리고 그들은 그 두 사람이 예수와 함께 다녔다는 사실을 알았지만,", "When they saw the courage of Peter and John and realized that they were unschooled, ordinary men, they were astonished and they took note that these men had been with Jesus."),
    ("사도행전 4:14 / Acts 4:14", "병 고침을 받은 사람이 그들 곁에 서 있는 것을 보고는, 아무 트집도 잡을 수 없었다.", "But since they could see the man who had been healed standing there with them, there was nothing they could say."),
    ("사도행전 4:15 / Acts 4:15", "그래서 그들은 그 두 사람에게 명령하여 의회에서 나가게 한 뒤에, 서로 의논하면서 말하였다.", "So they ordered them to withdraw from the Sanhedrin and then conferred together."),
    ("사도행전 4:16 / Acts 4:16", "\"이 사람들을 어떻게 하면 좋겠습니까? 그들로 말미암아 기적이 일어났다는 사실은, 예루살렘에 사는 모든 사람이 다 알고 있고, 우리도 이것을 부인할 수 없습니다.", "\"What are we going to do with these men?\" they asked. \"Everyone living in Jerusalem knows they have performed a notable sign, and we cannot deny it."),
    ("사도행전 4:17 / Acts 4:17", "다만 이 소문이 사람들에게 더 퍼지지 못하게, 앞으로는 이 이름으로 아무에게도 말하지 말라고, 그들에게 엄중히 경고합시다.\"", "But to stop this thing from spreading any further among the people, we must warn them to speak no longer to anyone in this name.\""),
    ("사도행전 4:18 / Acts 4:18", "그런 다음에, 그들은 그 두 사람을 불러서, 절대로 예수의 이름으로 말하지도 말고 가르치지도 말라고 명령하였다.", "Then they called them in again and commanded them not to speak or teach at all in the name of Jesus."),
    ("사도행전 4:19 / Acts 4:19", "그때에 베드로와 요한은 대답하였다. \"하나님의 말씀을 듣는 것보다, 당신들의 말을 듣는 것이, 하나님 보시기에 옳은 일인가를 판단해 보십시오.", "But Peter and John replied, \"Which is right in God's eyes: to listen to you, or to him? You be the judges!"),
    ("사도행전 4:20 / Acts 4:20", "우리는 보고 들은 것을 말하지 않을 수 없습니다.\"", "As for us, we cannot help speaking about what we have seen and heard.\""),
    ("사도행전 4:21 / Acts 4:21", "백성이 모두 그 일어난 일로 하나님께 영광을 돌리고 있으므로, 그들은 사도들을 처벌할 방도가 없어서, 다시 위협만 하고서 놓아 보냈다.", "After further threats they let them go. They could not decide how to punish them, because all the people were praising God for what had happened."),
    ("사도행전 4:22 / Acts 4:22", "이 기적으로 병이 나은 이는 마흔 살이 넘은 사람이다.", "For the man who was miraculously healed was over forty years old."),
    ("사도행전 4:23 / Acts 4:23", "베드로와 요한은 풀려나는 길로 동료들에게로 가서, 대제사장들과 장로들이 한 말을 낱낱이 일렀다.", "On their release, Peter and John went back to their own people and reported all that the chief priests and the elders had said to them."),
    ("사도행전 4:24 / Acts 4:24", "동료들은 이 말을 듣고서, 다같이 하나님께 부르짖어 아뢰었다. \"하늘과 땅과 바다와 그 안에 있는 모든 것을 지으신 주님,", "When they heard this, they raised their voices together in prayer to God. \"Sovereign Lord,\" they said, \"you made the heavens and the earth and the sea, and everything in them.\""),
    ("사도행전 4:25 / Acts 4:25", "주님께서는 주님의 종인 우리의 조상 다윗의 입을 빌어서, 성령으로 이렇게 말씀하셨습니다. '어찌하여 이방 민족이 날뛰며, 뭇 백성이 헛된 일을 꾀하였는가?", "You spoke by the Holy Spirit through the mouth of your servant, our father David: 'Why do the nations rage and the peoples plot in vain?"),
    ("사도행전 4:26 / Acts 4:26", "세상 임금들이 들고일어나고, 통치자들이 함께 모여서, 주님과 그의 메시아에게 대적하였다.'", "The kings of the earth rise up and the rulers band together against the Lord and against his anointed one.'"),
    ("사도행전 4:27 / Acts 4:27", "사실, 헤롯과 본디오 빌라도가 이방 사람들과 이스라엘 백성과 한패가 되어, 이 성에 모여서, 주님께서 기름 부으신 거룩한 종 예수를 대적하여,", "Indeed Herod and Pontius Pilate met together with the Gentiles and the people of Israel in this city to conspire against your holy servant Jesus, whom you anointed."),
    ("사도행전 4:28 / Acts 4:28", "주님의 권능과 뜻으로 미리 정하여 두신 일들을 모두 행하였습니다.", "They did what your power and will had decided beforehand should happen."),
    ("사도행전 4:29 / Acts 4:29", "주님, 이제 그들의 위협을 내려다보시고, 주님의 종들이 참으로 담대하게 주님의 말씀을 말할 수 있게 해주십시오.", "Now, Lord, consider their threats and enable your servants to speak your word with great boldness."),
    ("사도행전 4:30 / Acts 4:30", "그리고 주님께서 능력의 손을 뻗치시어 병을 낫게 해주시고, 주님의 거룩한 종 예수의 이름으로 표징과 놀라운 일들이 일어나게 해주십시오.\"", "Stretch out your hand to heal and perform signs and wonders through the name of your holy servant Jesus.\""),
    ("사도행전 4:31 / Acts 4:31", "그들이 기도를 마치니, 그들이 모여 있는 곳이 흔들리고, 그들은 모두 성령으로 충만해서, 하나님의 말씀을 담대히 말하게 되었다.", "After they prayed, the place where they were meeting was shaken. And they were all filled with the Holy Spirit and spoke the word of God boldly."),
    ("사도행전 4:32 / Acts 4:32", "많은 신도가 다 한 마음과 한 뜻이 되어서, 아무도 자기 소유를 자기 것이라고 하지 않고, 모든 것을 공동으로 사용하였다.", "All the believers were one in heart and mind. No one claimed that any of their possessions was their own, but they shared everything they had."),
    ("사도행전 4:33 / Acts 4:33", "사도들은 큰 능력으로 주 예수의 부활을 증언하였고, 사람들은 모두 큰 은혜를 받았다.", "With great power the apostles continued to testify to the resurrection of the Lord Jesus. And God's grace was so powerfully at work in them all"),
    ("사도행전 4:34 / Acts 4:34", "그들 가운데는 가난한 사람이 한 사람도 없었다. 땅이나 집을 가진 사람들은 그것을 팔아서, 그 판 돈을 가져다가", "that there were no needy persons among them. For from time to time those who owned land or houses sold them, brought the money from the sales"),
    ("사도행전 4:35 / Acts 4:35", "사도들의 발 앞에 놓았고, 사도들은 각 사람에게 필요에 따라 나누어주었다.", "and put it at the apostles' feet, and it was distributed to anyone who had need."),
    ("사도행전 4:36 / Acts 4:36", "키프로스 태생으로, 레위 사람이요, 사도들에게서 바나바 곧 '위로의 아들'이라는 뜻의 별명을 받은 요셉이,", "Joseph, a Levite from Cyprus, whom the apostles called Barnabas (which means \"son of encouragement\"),"),
    ("사도행전 4:37 / Acts 4:37", "자기가 가지고 있는 밭을 팔아서, 그 돈을 가져다가 사도들의 발 앞에 놓았다.", "sold a field he owned and brought the money and put it at the apostles' feet."),
]

DARK_NAVY  = RGBColor(0x0D, 0x1B, 0x3E)
GOLD       = RGBColor(0xD4, 0xAF, 0x37)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GOLD = RGBColor(0xF0, 0xD9, 0x80)  # slightly lighter for English text
DIVIDER    = RGBColor(0xD4, 0xAF, 0x37)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

def add_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, text, left, top, width, height,
                font_size=22, bold=False, color=WHITE,
                align=PP_ALIGN.CENTER, font_name=None):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    if font_name:
        run.font.name = font_name
    else:
        run.font.name = "Malgun Gothic" if any(ord(c) > 127 for c in text) else "Georgia"

def make_slide(prs, ref, korean, english):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, DARK_NAVY)
    W, H = SLIDE_W, SLIDE_H
    m = Inches(0.65)

    # Top gold bar
    b = slide.shapes.add_shape(1, Inches(0), Inches(0), W, Inches(0.07))
    b.fill.solid(); b.fill.fore_color.rgb = GOLD; b.line.fill.background()

    # Verse reference (centered, top)
    add_textbox(slide, ref, m, Inches(0.12), W - m*2, Inches(0.55),
                font_size=17, bold=True, color=GOLD, align=PP_ALIGN.CENTER,
                font_name="Malgun Gothic")

    # Korean text (top half)
    add_textbox(slide, korean, m, Inches(0.85), W - m*2, Inches(2.7),
                font_size=24, color=WHITE, align=PP_ALIGN.CENTER,
                font_name="Malgun Gothic")

    # Divider line
    div = slide.shapes.add_shape(1, Inches(1.5), Inches(3.7), W - Inches(3.0), Inches(0.04))
    div.fill.solid(); div.fill.fore_color.rgb = DIVIDER; div.line.fill.background()

    # English text (bottom half)
    add_textbox(slide, english, m, Inches(3.85), W - m*2, Inches(2.7),
                font_size=22, color=LIGHT_GOLD, align=PP_ALIGN.CENTER,
                font_name="Georgia")

    # Bottom gold bar
    b2 = slide.shapes.add_shape(1, Inches(0), H - Inches(0.07), W, Inches(0.07))
    b2.fill.solid(); b2.fill.fore_color.rgb = GOLD; b2.line.fill.background()

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

for (ref, korean, english) in verses:
    make_slide(prs, ref, korean, english)

prs.save("Acts4_Korean_English.pptx")
print("Done! 37 slides saved as Acts4_Korean_English.pptx")