import sys
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Usage ──────────────────────────────────────────────────────────────
# python make_bilingual_pptx.py korean.txt english.txt [output.pptx]
# ───────────────────────────────────────────────────────────────────────

REF_PATTERN = re.compile(r'^.{1,40}\s+\d+:\d+\s*$')

def extract_verse_num(ref):
    """Extract just the chapter:verse digits (e.g. '26:7') from a reference string."""
    m = re.search(r'(\d+:\d+)', ref)
    return m.group(1) if m else ref


def parse_verses(filepath):
    """
    Parse a scripture text file into a list of (reference, text) tuples.

    Format expected (Layout A — text before reference):
        Some verse text here...

        BookName 26:1

        Next verse text...

        BookName 26:2

    Handles split verses: if the same reference number appears multiple times
    (because a long verse was broken across lines in the source), all text
    fragments are joined into a single entry.

    Also handles trailing text after the final reference (appended to last verse).
    """
    with open(filepath, encoding='utf-8') as f:
        lines = [l.strip() for l in f.readlines()]

    # Drop blank lines
    lines = [l for l in lines if l]

    # --- Pass 1: collect (ref, text_before_ref) segments ---
    segments = []   # list of [ref_string, text_string]
    pending = []    # text lines accumulated before next ref

    for line in lines:
        if REF_PATTERN.match(line):
            segments.append([line.strip(), " ".join(pending).strip()])
            pending = []
        else:
            pending.append(line)

    # Trailing text after the last reference — append to last segment
    if pending and segments:
        trailing = " ".join(pending).strip()
        if trailing:
            segments[-1][1] = (segments[-1][1] + " " + trailing).strip()

    # Drop segments with empty text
    segments = [[r, t] for r, t in segments if t]

    # --- Pass 2: merge consecutive segments with the same verse number ---
    merged = []
    for ref, text in segments:
        num = extract_verse_num(ref)
        if merged and extract_verse_num(merged[-1][0]) == num:
            merged[-1][1] = merged[-1][1] + " " + text
        else:
            merged.append([ref, text])

    return [(ref, text.strip()) for ref, text in merged]


def match_verses(ko_verses, en_verses):
    """
    Pair Korean and English verses by their chapter:verse number.
    Returns list of (ref_ko, text_ko, ref_en, text_en).
    Prints a report of any unmatched verses.
    """
    ko_dict = {}
    for ref, text in ko_verses:
        num = extract_verse_num(ref)
        ko_dict[num] = (ref, text)

    en_dict = {}
    for ref, text in en_verses:
        num = extract_verse_num(ref)
        en_dict[num] = (ref, text)

    all_nums = sorted(ko_dict.keys() | en_dict.keys(),
                      key=lambda x: list(map(int, x.split(':'))))

    matched = []
    only_ko = []
    only_en = []

    for num in all_nums:
        if num in ko_dict and num in en_dict:
            ref_ko, text_ko = ko_dict[num]
            ref_en, text_en = en_dict[num]
            matched.append((ref_ko, text_ko, ref_en, text_en))
        elif num in ko_dict:
            only_ko.append(num)
        else:
            only_en.append(num)

    if only_ko:
        print(f"  ⚠  Korean-only verses (no English match): {only_ko}")
    if only_en:
        print(f"  ⚠  English-only verses (no Korean match): {only_en}")

    return matched


# ── Colors ─────────────────────────────────────────────────────────────
BLACK      = RGBColor(0x00, 0x00, 0x00)
GOLD       = RGBColor(0xD4, 0xAF, 0x37)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GOLD = RGBColor(0xF0, 0xD9, 0x80)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def add_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, text, left, top, width, height,
                font_size=22, bold=False, color=WHITE,
                align=PP_ALIGN.CENTER, font_name=None):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name or (
        "Malgun Gothic" if any(ord(c) > 127 for c in text) else "Georgia"
    )


def make_slide(prs, ref_ko, korean, ref_en, english):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_bg(slide, BLACK)
    W, H = SLIDE_W, SLIDE_H
    m = Inches(0.65)

    # Top gold bar
    b = slide.shapes.add_shape(1, Inches(0), Inches(0), W, Inches(0.07))
    b.fill.solid(); b.fill.fore_color.rgb = GOLD; b.line.fill.background()

    # Combined reference label
    ref_label = f"{ref_ko}  /  {ref_en}"
    add_textbox(slide, ref_label, m, Inches(0.15), W - m * 2, Inches(0.55),
                font_size=27, bold=True, color=GOLD, align=PP_ALIGN.CENTER,
                font_name="Malgun Gothic")

    # Korean text (top half)
    add_textbox(slide, korean, m, Inches(0.85), W - m * 2, Inches(2.9),
                font_size=36, color=WHITE, align=PP_ALIGN.CENTER,
                font_name="Malgun Gothic")

    # Divider
    div = slide.shapes.add_shape(1, Inches(1.5), Inches(3.85), W - Inches(3.0), Inches(0.04))
    div.fill.solid(); div.fill.fore_color.rgb = GOLD; div.line.fill.background()

    # English text (bottom half)
    add_textbox(slide, english, m, Inches(4.0), W - m * 2, Inches(2.9),
                font_size=36, color=LIGHT_GOLD, align=PP_ALIGN.CENTER,
                font_name="Georgia")

    # Bottom gold bar
    b2 = slide.shapes.add_shape(1, Inches(0), H - Inches(0.07), W, Inches(0.07))
    b2.fill.solid(); b2.fill.fore_color.rgb = GOLD; b2.line.fill.background()


def main():
    if len(sys.argv) < 3:
        print("Usage: python make_bilingual_pptx.py korean.txt english.txt [output.pptx]")
        sys.exit(1)

    ko_file  = sys.argv[1]
    en_file  = sys.argv[2]
    out_file = sys.argv[3] if len(sys.argv) > 3 else "bilingual_scripture.pptx"

    print(f"Parsing {ko_file} ...")
    ko_verses = parse_verses(ko_file)
    print(f"  → {len(ko_verses)} verses after merging")

    print(f"Parsing {en_file} ...")
    en_verses = parse_verses(en_file)
    print(f"  → {len(en_verses)} verses after merging")

    print("Matching verses by number...")
    pairs = match_verses(ko_verses, en_verses)
    print(f"  → {len(pairs)} matched pairs")

    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    for ref_ko, text_ko, ref_en, text_en in pairs:
        make_slide(prs, ref_ko, text_ko, ref_en, text_en)

    prs.save(out_file)
    print(f"✓ Saved {len(pairs)} slides → {out_file}")


if __name__ == "__main__":
    main()