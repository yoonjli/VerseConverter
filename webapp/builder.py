import io
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

BLACK      = RGBColor(0x00, 0x00, 0x00)
GOLD       = RGBColor(0xD4, 0xAF, 0x37)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GOLD = RGBColor(0xF0, 0xD9, 0x80)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def _add_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_textbox(slide, text, left, top, width, height,
                 font_size=22, bold=False, color=WHITE,
                 align=PP_ALIGN.CENTER, font_name=None):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name or (
        "Malgun Gothic" if any(ord(c) > 127 for c in text) else "Georgia"
    )


def _make_slide(prs, ref_ko, korean, ref_en, english):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, BLACK)
    W, H = SLIDE_W, SLIDE_H
    m = Inches(0.65)

    # Top gold bar
    b = slide.shapes.add_shape(1, Inches(0), Inches(0), W, Inches(0.07))
    b.fill.solid(); b.fill.fore_color.rgb = GOLD; b.line.fill.background()

    # Reference label
    _add_textbox(slide, f"{ref_ko}  /  {ref_en}",
                 m, Inches(0.15), W - m * 2, Inches(0.55),
                 font_size=30, bold=True, color=GOLD,
                 align=PP_ALIGN.CENTER, font_name="Malgun Gothic")

    # Korean text
    _add_textbox(slide, korean,
                 m, Inches(0.85), W - m * 2, Inches(2.9),
                 font_size=36, color=WHITE,
                 align=PP_ALIGN.CENTER, font_name="Malgun Gothic")

    # Divider
    div = slide.shapes.add_shape(1, Inches(1.5), Inches(3.85), W - Inches(3.0), Inches(0.04))
    div.fill.solid(); div.fill.fore_color.rgb = GOLD; div.line.fill.background()

    # English text
    _add_textbox(slide, english,
                 m, Inches(4.0), W - m * 2, Inches(2.9),
                 font_size=36, color=LIGHT_GOLD,
                 align=PP_ALIGN.CENTER, font_name="Georgia")

    # Bottom gold bar
    b2 = slide.shapes.add_shape(1, Inches(0), H - Inches(0.07), W, Inches(0.07))
    b2.fill.solid(); b2.fill.fore_color.rgb = GOLD; b2.line.fill.background()


def build_pptx(matched_verses: list[dict]) -> bytes:
    """
    Build a PPTX from a list of matched verse dicts.
    Returns the PPTX as raw bytes.
    """
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    for v in matched_verses:
        _make_slide(prs, v["ref_ko"], v["text_ko"], v["ref_en"], v["text_en"])

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()
